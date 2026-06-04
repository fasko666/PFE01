"""
Genere la presentation PowerPoint FreeNest — version maj. juin 2026
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ── Palette ──────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0,  52, 102)
BLUE_MAIN  = RGBColor(0,  82, 136)
BLUE_LIGHT = RGBColor(0, 120, 215)
ORANGE     = RGBColor(230,126,  34)
GREEN      = RGBColor(39, 174,  96)
RED        = RGBColor(192,  57,  43)
WHITE      = RGBColor(255,255,255)
GREY_LIGHT = RGBColor(245,247,250)
GREY_TEXT  = RGBColor(100,100,100)

W  = Inches(13.33)
H  = Inches(7.5)
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]

BASE    = r"C:\Users\Pro\Desktop\PFE O1\documentation"
UML_DIR = os.path.join(BASE, "uml_diagrams")

def add_rect(slide, x, y, w, h, fill=BLUE_MAIN, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_img(slide, path, x, y, w, h=None):
    if os.path.exists(path):
        if h:
            slide.shapes.add_picture(path, x, y, w, h)
        else:
            slide.shapes.add_picture(path, x, y, w)

def bullet_list(slide, items, x, y, w, h, size=14, color=BLUE_DARK, indent=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("    " if indent else "  ") + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)

def slide_bg(slide, color=GREY_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def slide_header(slide, title, subtitle=None, accent=BLUE_LIGHT):
    add_rect(slide, 0, 0, W, Inches(1.1), fill=BLUE_DARK)
    add_text(slide, title, Inches(0.4), Inches(0.12), Inches(10), Inches(0.7),
             size=26, bold=True, color=WHITE)
    add_rect(slide, 0, Inches(1.1), W, Inches(0.06), fill=accent)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(1.2), Inches(10), Inches(0.5),
                 size=14, color=GREY_TEXT)

def divider(slide, y, color=BLUE_LIGHT):
    add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.04), fill=color)

def stat_box(slide, label, value, x, y, w=Inches(2.8), h=Inches(1.3), bg=BLUE_MAIN):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, value, x, y+Inches(0.12), w, Inches(0.7),
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y+Inches(0.75), w, Inches(0.45),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)

def badge(slide, text, x, y, bg=ORANGE):
    add_rect(slide, x, y, Inches(2.1), Inches(0.38), fill=bg)
    add_text(slide, text, x+Inches(0.05), y+Inches(0.05), Inches(2.0), Inches(0.3),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def two_col(slide, left_items, right_items, y_start=Inches(2.0), size=13):
    bullet_list(slide, left_items,  Inches(0.5), y_start, Inches(6.0), Inches(4.5), size=size)
    bullet_list(slide, right_items, Inches(6.8), y_start, Inches(6.0), Inches(4.5), size=size)

# ════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITRE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BLUE_DARK)
add_rect(s, 0, Inches(2.8), W, Inches(2.2), fill=BLUE_MAIN)
add_text(s, "FreeNest", 0, Inches(0.6), W, Inches(1.2),
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Marketplace Freelance Full-Stack avec IA", 0, Inches(1.7), W, Inches(0.8),
         size=20, color=RGBColor(160,200,255), align=PP_ALIGN.CENTER)
add_text(s, "Laravel 12  |  React 19  |  Stripe  |  Docker  |  Ollama/Mistral 7B",
         0, Inches(2.9), W, Inches(0.6), size=15, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Projet de Fin d'Etudes — Ayoub Elmernissi — Juin 2026",
         0, Inches(3.6), W, Inches(0.5), size=13, color=RGBColor(180,220,255), align=PP_ALIGN.CENTER)
add_rect(s, Inches(4), Inches(4.5), Inches(5.33), Inches(0.05), fill=ORANGE)
add_text(s, "ayoubelmerniss55@gmail.com",
         0, Inches(4.8), W, Inches(0.5), size=13, color=RGBColor(200,200,200), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
#  SLIDE 2 — PLAN
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Plan de la Presentation")
plans = [
    ("01", "Contexte et Problematique"),
    ("02", "Presentation de FreeNest"),
    ("03", "Stack Technologique"),
    ("04", "Architecture et CI/CD"),
    ("05", "Authentification, 2FA et KYC"),
    ("06", "Marketplace et Catalogue"),
    ("07", "Contrats, Jalons et Temps"),
    ("08", "Paiements Escrow et Stripe"),
    ("09", "Messagerie Temps Reel et Push"),
    ("10", "Intelligence Artificielle"),
    ("11", "Agences et Gestion Talents"),
    ("12", "Administration et Finance"),
    ("13", "Base de Donnees (35+ tables)"),
    ("14", "Diagrammes UML"),
    ("15", "Deploiement Docker"),
    ("16", "Demo et Resultats"),
    ("17", "Bilan et Perspectives"),
]
for i, (num, title) in enumerate(plans):
    col = 0 if i < 9 else 1
    row = i if i < 9 else i - 9
    x = Inches(0.5 + col * 6.5)
    y = Inches(1.5 + row * 0.55)
    add_rect(s, x, y, Inches(0.45), Inches(0.38), fill=BLUE_LIGHT)
    add_text(s, num, x, y+Inches(0.04), Inches(0.45), Inches(0.32),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x+Inches(0.5), y+Inches(0.06), Inches(5.8), Inches(0.32),
             size=12, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 3 — CONTEXTE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Contexte et Problematique")
add_text(s, "Le marche mondial du freelance", Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         size=18, bold=True, color=BLUE_DARK)
stat_box(s, "Freelancers dans le monde", "1.57 Md", Inches(0.4), Inches(2.0), bg=BLUE_MAIN)
stat_box(s, "Marche 2025 (Mds $)",         "455 Md$",  Inches(3.4), Inches(2.0), bg=BLUE_LIGHT)
stat_box(s, "Croissance annuelle",           "+15%",     Inches(6.4), Inches(2.0), bg=GREEN)
stat_box(s, "Commission Upwork",             "20%+",     Inches(9.4), Inches(2.0), bg=RED)

add_text(s, "Problematique : Les plateformes existantes souffrent de :",
         Inches(0.5), Inches(3.5), Inches(12), Inches(0.4), size=14, bold=True, color=BLUE_DARK)
bullet_list(s, [
    "* Commissions elevees (20%+ chez Upwork)",
    "* Interfaces vieillissantes, UX mediocre",
    "* Manque de transparence sur les paiements",
    "* Aucune IA integree pour aider freelancers et clients",
    "* Pas d'agences, pas de catalogue de services flexibles",
], Inches(0.5), Inches(3.9), Inches(12), Inches(2.5), size=13, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 4 — PRESENTATION FREENEST
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Presentation de FreeNest", "Marketplace freelance full-stack de niveau production")

add_text(s, "FreeNest — Alternative moderne et complete aux grandes plateformes",
         Inches(0.5), Inches(1.3), Inches(12), Inches(0.5), size=14, color=GREY_TEXT)

features = [
    "Gestion de contrats et jalons complete",
    "Paiements escrow securises via Stripe",
    "Messagerie temps reel (WebSockets)",
    "Catalogue de services productises",
    "Gestion d'agences freelance",
    "IA locale (Ollama / Mistral 7B)",
    "KYC et authentification 2FA",
    "Notifications Push Web",
    "Deploiement Docker + CI/CD",
    "Centre fiscal (W-9, VAT, W-8BEN)",
]
for i, f in enumerate(features):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.9 + row * 0.8)
    add_rect(s, x, y, Inches(5.8), Inches(0.62), fill=WHITE)
    add_rect(s, x, y, Inches(0.08), Inches(0.62), fill=BLUE_LIGHT)
    add_text(s, "  " + f, x+Inches(0.15), y+Inches(0.12), Inches(5.5), Inches(0.4),
             size=13, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 5 — TECHNOLOGIES
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Stack Technologique")

add_text(s, "BACKEND", Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
         size=15, bold=True, color=BLUE_MAIN)
add_text(s, "FRONTEND", Inches(7.0), Inches(1.3), Inches(6), Inches(0.4),
         size=15, bold=True, color=ORANGE)

backend = [
    "Laravel 12 — API REST (100+ endpoints)",
    "PHP 8.2 — types stricts, attributs",
    "Sanctum — Bearer tokens",
    "Socialite — OAuth Google",
    "Stripe PHP SDK — paiements + webhooks",
    "Laravel Echo — broadcast WebSocket",
    "Spatie Permissions — roles et droits",
    "Queues + Scheduler — taches async",
]
frontend = [
    "React 19 — Concurrent Mode + Suspense",
    "Vite 6 — build ultra-rapide, HMR",
    "TailwindCSS 3.4 — dark/light mode",
    "Zustand 5.0 — state management",
    "Framer Motion 12 — animations",
    "Socket.IO 4.8 — temps reel",
    "Three.js + R3F — globe 3D",
    "Recharts 3.8 — graphiques",
]
bullet_list(s, backend,  Inches(0.5), Inches(1.7), Inches(6.0), Inches(4.5), size=13, color=BLUE_DARK)
bullet_list(s, frontend, Inches(7.0), Inches(1.7), Inches(6.0), Inches(4.5), size=13, color=BLUE_DARK)

divider(s, Inches(6.6))
infra = ["Docker + Compose", "GitHub Actions CI/CD", "MySQL 8.0 (35+ tables)", "Redis (cache/queues)", "Ollama / Mistral 7B (IA locale)", "Nginx (frontend + backend)"]
bullet_list(s, infra, Inches(0.5), Inches(6.65), Inches(12.5), Inches(0.8), size=12, color=GREY_TEXT)

# ════════════════════════════════════════════════════════════════
#  SLIDE 6 — ARCHITECTURE ET CI/CD
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Architecture et CI/CD")
add_text(s, "Separation stricte Frontend / Backend / IA — REST + WebSocket",
         Inches(0.5), Inches(1.2), Inches(12), Inches(0.4), size=13, color=GREY_TEXT)

layers = [
    ("React 19 SPA (Vite + Nginx)", BLUE_LIGHT, Inches(0.5), Inches(1.7)),
    ("Laravel 12 API REST",          BLUE_MAIN,  Inches(0.5), Inches(2.6)),
    ("MySQL 8.0 + Redis",            BLUE_DARK,  Inches(0.5), Inches(3.5)),
    ("Ollama / Mistral 7B",          GREEN,      Inches(7.0), Inches(2.6)),
    ("Stripe API + Webhooks",         ORANGE,     Inches(7.0), Inches(3.5)),
    ("Google OAuth 2.0",             RGBColor(66,133,244), Inches(7.0), Inches(1.7)),
]
for label, color, x, y in layers:
    add_rect(s, x, y, Inches(5.5), Inches(0.65), fill=color)
    add_text(s, label, x+Inches(0.2), y+Inches(0.15), Inches(5.2), Inches(0.4),
             size=14, bold=True, color=WHITE)

add_text(s, "CI/CD GitHub Actions", Inches(0.5), Inches(4.5), Inches(7), Inches(0.4),
         size=14, bold=True, color=BLUE_DARK)
bullet_list(s, [
    "backend.yml : PHP CS Fixer + PHPStan + PHPUnit + Docker build",
    "frontend.yml : npm ci + ESLint + Vite build + Vitest",
    "Declenchement sur chaque push/PR — echec bloque le merge",
], Inches(0.5), Inches(4.9), Inches(12), Inches(1.8), size=13, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 7 — AUTH, 2FA, KYC
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Authentification, Securite et 2FA + KYC")

add_text(s, "Authentification", Inches(0.5), Inches(1.3), Inches(4), Inches(0.4),
         size=14, bold=True, color=BLUE_MAIN)
bullet_list(s, [
    "Email + mot de passe (bcrypt)",
    "Google OAuth 2.0 (Socialite)",
    "Reset par email (2 etapes)",
    "Throttle: 10 req/min auth",
    "Throttle: 5 req/min reset",
], Inches(0.5), Inches(1.7), Inches(4.2), Inches(2.5), size=12, color=BLUE_DARK)

add_text(s, "2FA TOTP", Inches(4.9), Inches(1.3), Inches(4), Inches(0.4),
         size=14, bold=True, color=GREEN)
bullet_list(s, [
    "Google Authenticator",
    "QR Code genere (secret TOTP)",
    "Confirmation code 6 chiffres",
    "8 codes de secours jetables",
    "Regeneration des codes",
], Inches(4.9), Inches(1.7), Inches(4.0), Inches(2.5), size=12, color=BLUE_DARK)

add_text(s, "KYC - Identite", Inches(9.2), Inches(1.3), Inches(4), Inches(0.4),
         size=14, bold=True, color=ORANGE)
bullet_list(s, [
    "Upload passeport / CNI",
    "Selfie de verification",
    "Stockage local securise",
    "File d'attente admin",
    "Approbation / rejet",
], Inches(9.2), Inches(1.7), Inches(4.0), Inches(2.5), size=12, color=BLUE_DARK)

divider(s, Inches(4.3))
add_text(s, "Securite Globale", Inches(0.5), Inches(4.4), Inches(12), Inches(0.4),
         size=14, bold=True, color=BLUE_DARK)
bullet_list(s, [
    "Audit Logs : toute action sensible tracee (user, action, IP, metadata JSON)",
    "Rate Limiting : auth (10/min), IA (20/min), search (60-120/min)",
    "Validation stricte Form Requests Laravel + Middleware EnsurePlan",
], Inches(0.5), Inches(4.8), Inches(12.5), Inches(2.0), size=13, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 8 — MARKETPLACE ET CATALOGUE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Marketplace et Catalogue de Services")

add_text(s, "Offres d'emploi", Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
         size=14, bold=True, color=BLUE_MAIN)
bullet_list(s, [
    "Publication avec budget min/max, type (horaire/fixe)",
    "Statuts : draft -> open -> in_progress -> completed",
    "Index FULLTEXT MySQL pour recherche ultra-rapide",
    "Filtres : categorie, budget, type, date",
    "Sauvegarde des offres favorites",
], Inches(0.5), Inches(1.7), Inches(6.2), Inches(2.5), size=12, color=BLUE_DARK)

add_text(s, "Catalogue (Fiverr-like)", Inches(7.0), Inches(1.3), Inches(6), Inches(0.4),
         size=14, bold=True, color=ORANGE)
bullet_list(s, [
    "Services productises (3 tiers : Basic/Standard/Premium)",
    "Images multiples, descriptions detaillees",
    "Commandes : checkout, livraison, validation",
    "Avis post-commande (1-5 etoiles)",
    "Moderation admin (approuver / rejeter)",
    "Sauvegarde des services favoris",
], Inches(7.0), Inches(1.7), Inches(6.0), Inches(2.5), size=12, color=BLUE_DARK)

divider(s, Inches(4.3))
add_text(s, "Recherche Globale + Propositions", Inches(0.5), Inches(4.4), Inches(12), Inches(0.4),
         size=14, bold=True, color=BLUE_DARK)
bullet_list(s, [
    "GET /api/search : recherche unifiee (offres + freelancers + services + agences)",
    "GET /api/search/suggest : autocomplete temps reel (throttle 120/min) — resultats < 50ms",
    "Propositions : cover letter + bid + jalons JSON | Generee par IA | Flag is_ai_generated",
], Inches(0.5), Inches(4.8), Inches(12.5), Inches(2.0), size=13, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 9 — CONTRATS, JALONS, TEMPS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Contrats, Jalons et Suivi du Temps")

statuses = [
    ("active",    BLUE_LIGHT),
    ("paused",    GREY_TEXT),
    ("completed", GREEN),
    ("disputed",  ORANGE),
    ("cancelled", RED),
    ("archived",  BLUE_DARK),
]
for i, (status, color) in enumerate(statuses):
    add_rect(s, Inches(0.4 + i * 2.15), Inches(1.4), Inches(1.9), Inches(0.45), fill=color)
    add_text(s, status, Inches(0.4 + i * 2.15), Inches(1.47), Inches(1.9), Inches(0.35),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Fonctionnalites avancees :", Inches(0.5), Inches(2.0), Inches(12), Inches(0.4),
         size=13, bold=True, color=BLUE_DARK)
two_col(s, [
    "Fichiers livres (versioning)",
    "Suivi du temps (start/stop timer)",
    "Extensions de deadline",
    "Analytics du contrat",
    "Journal d'activite chronologique",
], [
    "Export PDF (contrat + litige)",
    "Archive des contrats termines",
    "Resolution de litiges par admin",
    "Conversation liee au contrat",
    "Notifications temps reel",
], Inches(2.4), size=13)

divider(s, Inches(4.9))
add_text(s, "Workflow des Jalons :", Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
         size=13, bold=True, color=BLUE_DARK)
milestone_flow = ["pending", "->", "funded", "->", "submitted", "->", "approved"]
colors_flow = [GREY_TEXT, BLUE_DARK, ORANGE, BLUE_DARK, BLUE_LIGHT, BLUE_DARK, GREEN]
for i, (step, color) in enumerate(zip(milestone_flow, colors_flow)):
    add_text(s, step, Inches(0.5 + i * 1.7), Inches(5.4), Inches(1.5), Inches(0.5),
             size=13, bold=(step != "->"), color=color)

bullet_list(s, [
    "Jalon rejected -> retour a 'funded' | Jalon disputed -> admin arbitre | Factures hebdomadaires (WeeklyInvoice)",
], Inches(0.5), Inches(5.9), Inches(12.5), Inches(0.8), size=12, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 10 — PAIEMENTS ESCROW ET STRIPE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Systeme de Paiement Escrow et Stripe")

phases = [
    ("1. Depot", "Client paye via\nStripe Checkout\n(3D Secure)", BLUE_LIGHT),
    ("2. Escrow", "Client bloque\ndes fonds pour\nun jalon", ORANGE),
    ("3. Livraison", "Freelancer soumet\nson travail\n+ fichiers", GREEN),
    ("4. Liberation", "Client approuve\n-> 90% freelancer\n10% commission", BLUE_MAIN),
]
for i, (title, desc, color) in enumerate(phases):
    x = Inches(0.4 + i * 3.2)
    add_rect(s, x, Inches(1.5), Inches(2.9), Inches(1.2), fill=color)
    add_text(s, title, x, Inches(1.55), Inches(2.9), Inches(0.45),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x, Inches(2.0), Inches(2.9), Inches(0.75),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, "->", Inches(3.3 + i * 3.2), Inches(1.9), Inches(0.2), Inches(0.4),
                 size=20, bold=True, color=BLUE_DARK)

add_text(s, "Portefeuille a 3 balances :", Inches(0.5), Inches(3.1), Inches(12), Inches(0.4),
         size=14, bold=True, color=BLUE_DARK)
balances = [
    ("balance", "Fonds disponibles et retirables", GREEN),
    ("pending_balance", "Fonds en transit", ORANGE),
    ("escrow_balance", "Fonds bloques en garantie", BLUE_MAIN),
]
for i, (name, desc, color) in enumerate(balances):
    x = Inches(0.5 + i * 4.2)
    add_rect(s, x, Inches(3.55), Inches(3.8), Inches(0.8), fill=color)
    add_text(s, name, x+Inches(0.1), Inches(3.58), Inches(3.6), Inches(0.35),
             size=13, bold=True, color=WHITE)
    add_text(s, desc, x+Inches(0.1), Inches(3.93), Inches(3.6), Inches(0.35),
             size=11, color=WHITE)

divider(s, Inches(4.5))
add_text(s, "Stripe Connect + Retraits + Abonnements", Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
         size=14, bold=True, color=BLUE_DARK)
bullet_list(s, [
    "Stripe Connect Express : freelancers connectent leur IBAN pour virements",
    "Retraits soumis a approbation admin (KYC requis pour deblocage)",
    "Abonnements : Stripe Checkout + portail de facturation + invoices",
    "Webhooks : checkout.session.completed, account.updated, invoice.paid",
], Inches(0.5), Inches(5.0), Inches(12.5), Inches(2.2), size=12, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 11 — MESSAGERIE TEMPS REEL ET PUSH
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Messagerie Temps Reel et Notifications Push")

add_text(s, "Laravel Echo + Socket.IO 4.8 — 11 evenements WebSocket diffuses en temps reel",
         Inches(0.5), Inches(1.3), Inches(12), Inches(0.4), size=13, color=GREY_TEXT)

features_chat = [
    ("Envoyer / Recevoir",  "Texte, fichiers, images en temps reel"),
    ("Editer message",      "Historique des modifications preserve"),
    ("Supprimer message",   "Suppression douce (deleted_at)"),
    ("Reactions emoji",     "Toggle par message — MessageReaction"),
    ("Accusee de lecture",  "Marquer toute la conversation lue"),
    ("Indicateur frappe",   "UserTyping broadcast — 'X ecrit...'"),
    ("Livraison",           "MessageDelivered confirmation"),
    ("Recherche",           "Par conversation ou contenu message"),
    ("Piece jointe",        "Upload et partage de fichiers"),
    ("Lien contrat",        "Conversation liee au contrat actif"),
]
for i, (feat, desc) in enumerate(features_chat):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = Inches(0.4 + col * 6.5)
    y = Inches(1.8 + row * 0.82)
    add_rect(s, x, y, Inches(5.9), Inches(0.65), fill=WHITE)
    add_rect(s, x, y, Inches(0.08), Inches(0.65), fill=BLUE_LIGHT)
    add_text(s, feat, x+Inches(0.15), y+Inches(0.05), Inches(2.5), Inches(0.3),
             size=12, bold=True, color=BLUE_DARK)
    add_text(s, desc, x+Inches(0.15), y+Inches(0.33), Inches(5.6), Inches(0.28),
             size=11, color=GREY_TEXT)

divider(s, Inches(6.5))
add_text(s, "Push Notifications (VAPID) : Service Worker sw.js | subscribe / unsubscribe | Son (sound.js) + icone navigateur",
         Inches(0.5), Inches(6.55), Inches(12.5), Inches(0.7), size=12, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 12 — INTELLIGENCE ARTIFICIELLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Intelligence Artificielle — Ollama / Mistral 7B")

add_text(s, "IA locale — Aucun envoi vers OpenAI — Confidentialite totale — Fallback si indisponible",
         Inches(0.5), Inches(1.2), Inches(12), Inches(0.4), size=13, color=GREY_TEXT)

ai_features = [
    ("Generation de proposition", "Redige une proposition a partir du profil freelancer et de l'offre", BLUE_LIGHT),
    ("Matching freelancers",       "Recommande les meilleurs freelancers pour une offre client", GREEN),
    ("Analyse de profil",          "Conseils personnalises pour optimiser le profil freelancer", ORANGE),
    ("Chat assistant",             "Assistant IA conversationnel — questions/reponses contextuelles", BLUE_MAIN),
    ("Recherche intelligente",     "Recherche semantique au-dela des mots-cles exacts", RGBColor(142,68,173)),
]
for i, (title, desc, color) in enumerate(ai_features):
    y = Inches(1.7 + i * 0.92)
    add_rect(s, Inches(0.4), y, Inches(0.5), Inches(0.72), fill=color)
    add_text(s, str(i+1), Inches(0.4), y+Inches(0.15), Inches(0.5), Inches(0.4),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.0), y, Inches(11.9), Inches(0.72), fill=WHITE)
    add_text(s, title, Inches(1.1), y+Inches(0.05), Inches(4.5), Inches(0.35),
             size=13, bold=True, color=color)
    add_text(s, desc, Inches(1.1), y+Inches(0.38), Inches(11.7), Inches(0.3),
             size=11, color=GREY_TEXT)

divider(s, Inches(6.5))
bullet_list(s, [
    "Rate limit: 20 req/min | Logs tokens dans ai_histories | Fallback pre-genere si Ollama indisponible",
], Inches(0.5), Inches(6.55), Inches(12.5), Inches(0.7), size=12, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 13 — AGENCES ET TALENTS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Agences et Gestion des Talents")

add_text(s, "Agences", Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
         size=16, bold=True, color=BLUE_MAIN)
bullet_list(s, [
    "Creer une agence (logo, nom, description)",
    "Inviter des membres par email (token unique)",
    "Roles : owner / admin / member",
    "Transfert de propriete de l'agence",
    "Page publique dans le repertoire agences",
    "Accepter / decliner une invitation",
], Inches(0.5), Inches(1.75), Inches(6.0), Inches(3.0), size=13, color=BLUE_DARK)

add_text(s, "Gestion des Talents (Clients)", Inches(7.0), Inches(1.3), Inches(6), Inches(0.4),
         size=16, bold=True, color=ORANGE)
bullet_list(s, [
    "Saved Freelancers : signets rapides (un clic)",
    "Talent Lists : listes curees et nommees",
    "Ajouter / retirer des membres de liste",
    "Notes personnelles sur les freelancers",
    "Organisation par specialite ou projet",
], Inches(7.0), Inches(1.75), Inches(6.0), Inches(3.0), size=13, color=BLUE_DARK)

divider(s, Inches(4.9))
add_text(s, "Centre Fiscal", Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
         size=16, bold=True, color=GREEN)
bullet_list(s, [
    "Soumission : W-9 (USA), W-8BEN (international), VAT (Europe)",
    "Export PDF des documents fiscaux | Validation admin (approuver / rejeter avec motif)",
], Inches(0.5), Inches(5.4), Inches(12.5), Inches(1.5), size=13, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 14 — ADMINISTRATION ET FINANCE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Administration et Finance Admin")

admin_sections = [
    ("Dashboard Admin",     ["Stats : GMV, commissions, contrats", "Gestion utilisateurs (ban, verify)", "Analytics avancees"], BLUE_MAIN),
    ("Finance Admin",       ["Vue d'ensemble revenue plateforme", "File retraits (approuver / rejeter)", "Parametres : commission, limites"], ORANGE),
    ("KYC Admin",           ["File verifications en attente", "Acces securise aux documents", "Approuver / Rejeter + audit log"], GREEN),
    ("Moderation Catalogue",["Services en attente de validation", "Approuver / Rejeter publication", "Documents fiscaux W-9/W-8BEN/VAT"], RED),
]
for i, (title, items, color) in enumerate(admin_sections):
    col = i % 2
    row = i // 2
    x = Inches(0.4 + col * 6.5)
    y = Inches(1.5 + row * 2.5)
    add_rect(s, x, y, Inches(5.9), Inches(0.5), fill=color)
    add_text(s, title, x+Inches(0.2), y+Inches(0.08), Inches(5.6), Inches(0.38),
             size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text(s, "  > " + item, x+Inches(0.1), y+Inches(0.5+j*0.5), Inches(5.7), Inches(0.45),
                 size=12, color=BLUE_DARK)

divider(s, Inches(6.65))
add_text(s, "Middleware 'admin' + double verification controller — Audit logs sur toutes les actions sensibles",
         Inches(0.5), Inches(6.68), Inches(12.5), Inches(0.6), size=11, color=GREY_TEXT)

# ════════════════════════════════════════════════════════════════
#  SLIDE 15 — BASE DE DONNEES
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Base de Donnees — 35+ Tables MySQL")
add_img(s, os.path.join(UML_DIR, "06_erd.png"),
        Inches(0.3), Inches(1.3), Inches(13), Inches(5.9))

# ════════════════════════════════════════════════════════════════
#  SLIDE 16 — DIAGRAMMES UML
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Diagrammes UML — Vue d'Ensemble (20 diagrammes)")
uml_imgs = [
    ("01_use_case_global.png",  "Use Case Global"),
    ("02_use_case_freelancer.png", "Use Case Freelancer"),
    ("05_class_diagram.png",    "Diagramme de Classes"),
    ("14_activity_cycle.png",   "Activite Mission"),
    ("16_state_contrat.png",    "Etats Contrat"),
    ("17_state_milestone.png",  "Etats Jalon"),
]
positions = [
    (Inches(0.2),  Inches(1.4), Inches(4.2)),
    (Inches(4.6),  Inches(1.4), Inches(4.2)),
    (Inches(9.0),  Inches(1.4), Inches(4.2)),
    (Inches(0.2),  Inches(4.3), Inches(4.2)),
    (Inches(4.6),  Inches(4.3), Inches(4.2)),
    (Inches(9.0),  Inches(4.3), Inches(4.2)),
]
for (img, label), (x, y, w) in zip(uml_imgs, positions):
    path = os.path.join(UML_DIR, img)
    add_img(s, path, x, y, w, Inches(2.8))
    add_text(s, label, x, y+Inches(2.82), w, Inches(0.3),
             size=10, color=GREY_TEXT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
#  SLIDE 17 — DEPLOIEMENT DOCKER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Deploiement Docker et Infrastructure")

services = [
    ("Frontend\nReact 19 + Nginx", BLUE_LIGHT,  Inches(0.4)),
    ("Backend\nLaravel 12 API",    BLUE_MAIN,   Inches(3.3)),
    ("MySQL 8.0\n35+ tables",       ORANGE,      Inches(6.2)),
    ("Redis\nCache + Queues",       GREEN,        Inches(9.1)),
]
for label, color, x in services:
    add_rect(s, x, Inches(1.5), Inches(2.7), Inches(1.5), fill=color)
    add_text(s, label, x, Inches(1.65), Inches(2.7), Inches(1.2),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

services2 = [
    ("Soketi\nWebSocket",     RGBColor(142,68,173), Inches(0.4)),
    ("Ollama\nMistral 7B",    RGBColor(52,152,219), Inches(3.3)),
    ("Stripe\nPaiements",     RGBColor(99,91,255),  Inches(6.2)),
    ("GitHub Actions\nCI/CD", RGBColor(36,41,46),   Inches(9.1)),
]
for label, color, x in services2:
    add_rect(s, x, Inches(3.2), Inches(2.7), Inches(1.3), fill=color)
    add_text(s, label, x, Inches(3.35), Inches(2.7), Inches(1.0),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider(s, Inches(4.7))
add_text(s, "Configuration Production", Inches(0.5), Inches(4.8), Inches(12), Inches(0.4),
         size=14, bold=True, color=BLUE_DARK)
bullet_list(s, [
    "Dockerfile multi-stage backend (PHP-FPM + Nginx) + Dockerfile Nginx frontend",
    "nginx.conf : Gzip, Cache-Control (assets 1 an), SPA try_files, Service Worker",
    "docker-compose.yml : orchestration complete (6+ services)",
    "CI/CD : backend.yml + frontend.yml — echec bloque le merge automatiquement",
], Inches(0.5), Inches(5.2), Inches(12.5), Inches(2.0), size=13, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 18 — DEMO ET RESULTATS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Demo et Resultats")

stat_box(s, "Tables MySQL",      "35+",  Inches(0.3), Inches(1.5), bg=BLUE_MAIN)
stat_box(s, "Endpoints API",     "100+", Inches(3.4), Inches(1.5), bg=BLUE_LIGHT)
stat_box(s, "Pages React",       "50+",  Inches(6.5), Inches(1.5), bg=GREEN)
stat_box(s, "Composants React",  "80+",  Inches(9.6), Inches(1.5), bg=ORANGE)

stat_box(s, "Evenements WS",     "11",   Inches(0.3), Inches(3.0), bg=RGBColor(142,68,173))
stat_box(s, "Fonctionnalites IA","5",    Inches(3.4), Inches(3.0), bg=RGBColor(52,152,219))
stat_box(s, "Services Docker",   "6+",   Inches(6.5), Inches(3.0), bg=RED)
stat_box(s, "LOC (lignes code)", "25K+", Inches(9.6), Inches(3.0), bg=BLUE_DARK)

divider(s, Inches(4.5))
add_text(s, "Parcours utilisateur demontrable de bout en bout :", Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
         size=14, bold=True, color=BLUE_DARK)
bullet_list(s, [
    "Inscription -> 2FA -> KYC -> Publication offre -> Proposition IA -> Contrat -> Escrow Stripe -> Livraison -> Paiement",
    "Catalogue : service -> commande -> livraison -> avis  |  Agence : creation -> invitation -> collaboration",
    "Chat temps reel + reactions emoji + Push Web  |  Admin : finance, KYC, moderation catalogue",
], Inches(0.5), Inches(5.0), Inches(12.5), Inches(1.8), size=13, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 19 — BILAN
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Bilan et Perspectives")

add_text(s, "Competences acquises :", Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
         size=15, bold=True, color=BLUE_MAIN)
bullet_list(s, [
    "Architecture full-stack decouplee (API REST + SPA)",
    "Paiements complexes (Stripe, escrow, webhooks)",
    "Temps reel (Laravel Broadcasting + Socket.IO)",
    "Securite avancee (2FA, KYC, audit logs)",
    "DevOps complet (Docker, Nginx, CI/CD)",
    "Integration IA locale (Ollama, Mistral 7B)",
    "BDD relationnelle complexe (35+ tables)",
], Inches(0.5), Inches(1.75), Inches(6.0), Inches(3.5), size=13, color=BLUE_DARK)

add_text(s, "Perspectives :", Inches(7.0), Inches(1.3), Inches(6), Inches(0.4),
         size=15, bold=True, color=ORANGE)
bullet_list(s, [
    "Application mobile React Native",
    "Microservices (Auth / Payment / AI)",
    "IA multimodale (images, videos)",
    "Video consultations integrees",
    "Multi-devises, i18n international",
    "Kubernetes (orchestration prod)",
    "Certifications freelancers en ligne",
], Inches(7.0), Inches(1.75), Inches(6.0), Inches(3.5), size=13, color=BLUE_DARK)

divider(s, Inches(5.5))
add_text(s, "Difficultes surmontees :", Inches(0.5), Inches(5.6), Inches(12), Inches(0.4),
         size=14, bold=True, color=BLUE_DARK)
bullet_list(s, [
    "Transactions atomiques escrow (race conditions evitees avec BEGIN/COMMIT)",
    "Securite Stripe Webhooks (verification HMAC signature, idempotence des evenements)",
    "WebSockets en production (Soketi containerise, canaux prives Sanctum)",
], Inches(0.5), Inches(6.0), Inches(12.5), Inches(1.2), size=12, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
#  SLIDE 20 — CONCLUSION
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BLUE_DARK)
add_rect(s, 0, Inches(2.5), W, Inches(2.8), fill=BLUE_MAIN)
add_text(s, "Merci pour votre attention", 0, Inches(0.8), W, Inches(1.2),
         size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "FreeNest — Marketplace Freelance Full-Stack de Niveau Production",
         0, Inches(2.6), W, Inches(0.7), size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Laravel 12  |  React 19  |  Stripe  |  Docker  |  Ollama  |  CI/CD GitHub Actions",
         0, Inches(3.3), W, Inches(0.6), size=14, color=RGBColor(160,200,255), align=PP_ALIGN.CENTER)
add_rect(s, Inches(4), Inches(4.1), Inches(5.33), Inches(0.06), fill=ORANGE)
add_text(s, "Ayoub Elmernissi — ayoubelmerniss55@gmail.com — Juin 2026",
         0, Inches(4.3), W, Inches(0.5), size=14, color=RGBColor(200,220,255), align=PP_ALIGN.CENTER)
add_text(s, "Questions ?", 0, Inches(5.2), W, Inches(1.2),
         size=42, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Sauvegarde
out = r"C:\Users\Pro\Desktop\PFE O1\documentation\FreeNest_Presentation.pptx"
prs.save(out)
print("OK - Presentation PowerPoint creee : " + out)
print("Slides : " + str(len(prs.slides)))
